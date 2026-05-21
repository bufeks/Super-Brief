#!/usr/bin/env python3
"""Extract plain text from a brief file (PDF / PPTX / DOCX / TXT / MD).

Usage:
    python3 scripts/extract.py /absolute/path/to/brief.pdf

Output is written to stdout as plain text. Errors go to stderr.

Designed for use by the `research-brief` skill — the extracted text is
fed to Claude for decoding. Keep the output structure faithful to the
source (slides keep slide boundaries, tables keep row separators).
"""
from __future__ import annotations

import sys
from pathlib import Path


def extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        parts.append(f"--- Page {i} ---\n{text}")
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(f"[Notes] {notes}")
    return "\n".join(parts)


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def main() -> int:
    if len(sys.argv) < 2:
        print("Usage: extract.py <file>", file=sys.stderr)
        return 1

    path = Path(sys.argv[1])
    if not path.exists():
        print(f"File not found: {path}", file=sys.stderr)
        return 1

    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            text = extract_pdf(path)
        elif ext == ".pptx":
            text = extract_pptx(path)
        elif ext == ".docx":
            text = extract_docx(path)
        elif ext in (".txt", ".md"):
            text = path.read_text(encoding="utf-8")
        else:
            print(
                f"Unsupported file type: {ext}. "
                "Supported: .pdf .pptx .docx .txt .md",
                file=sys.stderr,
            )
            return 2
    except ImportError as e:
        print(
            f"Missing dependency for {ext}: {e}\n"
            "Install requirements: pip install -r requirements.txt",
            file=sys.stderr,
        )
        return 3

    text = text.strip()
    if not text:
        print("Extraction produced empty text — file may be image-based "
              "(OCR required).", file=sys.stderr)
        return 4

    print(text)
    return 0


if __name__ == "__main__":
    sys.exit(main())
